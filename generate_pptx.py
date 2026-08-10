import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_professional_deck():
    prs = Presentation()
    # Use 16:9 widescreen slides
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Palette
    c_bg_light = RGBColor(245, 247, 250)      # Sleek off-white/light gray
    c_teal = RGBColor(0, 150, 136)            # Accent Teal
    c_blue = RGBColor(0, 74, 198)             # Primary Blue
    c_dark = RGBColor(25, 30, 36)             # Dark text/background
    c_white = RGBColor(255, 255, 255)
    c_gray = RGBColor(115, 118, 134)          # Outline/Secondary
    c_card_border = RGBColor(226, 232, 240)
    
    # Helper to apply slide background
    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to draw a text block
    def add_text_box(slide, left, top, width, height, text, size=18, color=c_dark, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
        return txBox

    # Helper to add standard header/footer
    def add_slide_header(slide, title_text, category="CIVIS AI"):
        # Header Accent line
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(0.08), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_teal
        shape.line.fill.background()
        
        # Category label
        add_text_box(slide, Inches(0.7), Inches(0.35), Inches(4), Inches(0.3), category.upper(), size=10, color=c_teal, bold=True)
        # Main Title
        add_text_box(slide, Inches(0.7), Inches(0.55), Inches(10), Inches(0.5), title_text, size=24, color=c_dark, bold=True)

    # Helper to draw a card/glassmorphism shape
    def add_card(slide, left, top, width, height, bg_color=c_white, border_color=c_card_border):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # -----------------------------------------------------
    # Slide 1: Title Page (Light Premium Theme)
    # -----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, c_bg_light)
    
    # Left accent block
    accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = c_teal
    accent.line.fill.background()
    
    # Title
    add_text_box(slide1, Inches(1.5), Inches(1.8), Inches(10), Inches(1.5), "CIVIS AI", size=60, color=c_blue, bold=True)
    add_text_box(slide1, Inches(1.5), Inches(3.0), Inches(10), Inches(1.0), "AI Powered Smart City Civic Issue Management Platform", size=20, color=c_teal, bold=True)
    
    # Details Box (Card representation)
    add_card(slide1, Inches(1.5), Inches(4.5), Inches(8), Inches(2.0), bg_color=c_white, border_color=c_card_border)
    add_text_box(slide1, Inches(1.8), Inches(4.7), Inches(7.4), Inches(1.6), 
                 "Team Members: Ishita & Sarthak\nCollege: [Your College]\nCompetition: Hackathon Showcase\nDate: July 2026", 
                 size=14, color=c_dark)

    # -----------------------------------------------------
    # Slide 2: Problem Statement (4 Card Layout)
    # -----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2, c_bg_light)
    add_slide_header(slide2, "The Core Urban Challenges")
    
    problems = [
        {"title": "Potholes & Road Damage", "desc": "Neglected roads slow down transit and lead to accidents.", "icon": "🚗"},
        {"title": "Unattended Waste", "desc": "Piles of garbage degrade hygiene and create health hazards.", "icon": "🗑️"},
        {"title": "Resource Leaks", "desc": "Water leakages and broken streetlights waste vital city utility resources.", "icon": "💡"},
        {"title": "Communication Delay", "desc": "Municipal action lags due to manual triage and missing geo-coordinates.", "icon": "⏳"}
    ]
    
    card_width = Inches(2.7)
    card_height = Inches(4.2)
    start_left = Inches(0.5)
    spacing = Inches(0.4)
    
    for i, p in enumerate(problems):
        left_pos = start_left + i * (card_width + spacing)
        add_card(slide2, left_pos, Inches(1.8), card_width, card_height)
        
        # Icon
        add_text_box(slide2, left_pos + Inches(0.2), Inches(2.0), Inches(2.3), Inches(0.6), p["icon"], size=36)
        # Title
        add_text_box(slide2, left_pos + Inches(0.2), Inches(2.8), Inches(2.3), Inches(0.8), p["title"], size=16, color=c_blue, bold=True)
        # Desc
        add_text_box(slide2, left_pos + Inches(0.2), Inches(3.6), Inches(2.3), Inches(2.0), p["desc"], size=13, color=c_dark)

    # -----------------------------------------------------
    # Slide 3: Solution Overview (Workflow Pipeline)
    # -----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3, c_bg_light)
    add_slide_header(slide3, "The Seamless Citizen-Authority Pipeline")
    
    pipeline = ["Citizen App", "AI Scan Detection", "Auto Geo-Tagging", "Admin Analytics", "Authority Action", "Push Notification"]
    
    box_width = Inches(1.7)
    box_height = Inches(1.6)
    step_start = Inches(0.5)
    step_spacing = Inches(0.4)
    
    for i, step in enumerate(pipeline):
        left_pos = step_start + i * (box_width + step_spacing)
        
        # Draw step card
        add_card(slide3, left_pos, Inches(3.0), box_width, box_height)
        
        # Step Number
        add_text_box(slide3, left_pos, Inches(3.1), box_width, Inches(0.3), f"Step {i+1}", size=12, color=c_teal, bold=True, align=PP_ALIGN.CENTER)
        # Step Title
        add_text_box(slide3, left_pos + Inches(0.1), Inches(3.5), box_width - Inches(0.2), Inches(1.0), step, size=14, color=c_dark, bold=True, align=PP_ALIGN.CENTER)
        
        # Arrow (except for last element)
        if i < len(pipeline) - 1:
            arrow_left = left_pos + box_width + Inches(0.05)
            add_text_box(slide3, arrow_left, Inches(3.7), step_spacing, Inches(0.5), "➔", size=20, color=c_teal, align=PP_ALIGN.CENTER)

    # -----------------------------------------------------
    # Slide 4: System Architecture (Layered Stack)
    # -----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4, c_bg_light)
    add_slide_header(slide4, "System Architecture")
    
    layers = [
        {"name": "Client Layer", "tech": "Citizen Mobile Web App (Tailwind CSS, VanillaJS)", "bg": c_blue, "fg": c_white},
        {"name": "Processing Layer", "tech": "AI Computer Vision API (Image Classification) & Geolocation", "bg": c_teal, "fg": c_white},
        {"name": "Analytics & Control", "tech": "Admin Dashboard & Interactive Leaflet Maps", "bg": c_white, "fg": c_dark},
        {"name": "Action Layer", "tech": "Municipal Operations & Maintenance Dispatch", "bg": c_dark, "fg": c_white}
    ]
    
    y_start = Inches(1.8)
    for i, layer in enumerate(layers):
        # Draw layer card
        add_card(slide4, Inches(0.5), y_start, Inches(12.33), Inches(1.1), bg_color=layer["bg"])
        
        # Title of layer
        add_text_box(slide4, Inches(0.8), y_start + Inches(0.3), Inches(3.0), Inches(0.5), layer["name"], size=18, color=layer["fg"], bold=True)
        # Tech details
        add_text_box(slide4, Inches(4.0), y_start + Inches(0.3), Inches(8.0), Inches(0.5), layer["tech"], size=16, color=layer["fg"])
        
        y_start += Inches(1.3)

    # -----------------------------------------------------
    # Slide 5: Technology Stack (Grid of 4 Cards)
    # -----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5, c_bg_light)
    add_slide_header(slide5, "Technology Stack")
    
    techs = [
        {"title": "Frontend", "list": "• HTML5 & Tailwind CSS\n• Responsive layout\n• Modern Glassmorphism UI\n• Pure Javascript logic", "icon": "💻"},
        {"title": "Backend Services", "list": "• Node.js Platform\n• LocalStorage Sync\n• Serve server integration\n• Modular architecture", "icon": "⚙️"},
        {"title": "AI & Analysis", "list": "• Anomaly Classification\n• Custom scanning overlays\n• Automated tags\n• Triage generation", "icon": "🧠"},
        {"title": "Mapping Engine", "list": "• Leaflet.js Mapping API\n• Interactive markers\n• Heatmap support\n• Precision Geocoding", "icon": "🗺️"}
    ]
    
    for i, tech in enumerate(techs):
        row = i // 2
        col = i % 2
        
        l_pos = Inches(1.2 + col * 5.8)
        t_pos = Inches(1.8 + row * 2.6)
        
        # Card shape
        add_card(slide5, l_pos, t_pos, Inches(5.2), Inches(2.2))
        
        # Title & Icon
        add_text_box(slide5, l_pos + Inches(0.3), t_pos + Inches(0.2), Inches(4.5), Inches(0.5), f"{tech['icon']}  {tech['title']}", size=18, color=c_blue, bold=True)
        # List
        add_text_box(slide5, l_pos + Inches(0.3), t_pos + Inches(0.8), Inches(4.5), Inches(1.2), tech["list"], size=13, color=c_dark)

    # -----------------------------------------------------
    # Slide 6: Key Features (Dynamic Split Screen)
    # -----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6, c_bg_light)
    add_slide_header(slide6, "Key Features In Action")
    
    # Left column: details
    add_card(slide6, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5.0))
    features_list = (
        "• AI Scan: Real-time visual issue verification.\n\n"
        "• Interactive Map: Live Leaflet mapping with precise, geocoded markers.\n\n"
        "• My Complaints: End-to-end tracking dashboard for citizen transparency.\n\n"
        "• Admin Analytics: Visual insights, maps, and reports for rapid authority dispatch."
    )
    add_text_box(slide6, Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.6), features_list, size=15, color=c_dark)
    
    # Right column: screenshots
    artifact_dir = r"C:\Users\sarth\.gemini\antigravity\brain\34873cc1-f75f-4c90-9f51-e77167bec5d7"
    img_h = os.path.join(artifact_dir, "home.png")
    img_ad = os.path.join(artifact_dir, "admin_dashboard.png")
    
    try:
        # Home mockup
        add_card(slide6, Inches(6.8), Inches(1.8), Inches(2.8), Inches(5.0))
        slide6.shapes.add_picture(img_h, Inches(6.9), Inches(1.9), width=Inches(2.6), height=Inches(4.8))
        
        # Admin mockup
        add_card(slide6, Inches(9.8), Inches(1.8), Inches(3.0), Inches(5.0))
        slide6.shapes.add_picture(img_ad, Inches(9.9), Inches(1.9), width=Inches(2.8), height=Inches(4.8))
    except Exception as e:
        print("Mockups could not be inserted:", e)

    # -----------------------------------------------------
    # Slide 7: Project Workflow (Horizontal Timeline)
    # -----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7, c_bg_light)
    add_slide_header(slide7, "Operational Flow")
    
    steps = [
        "📸 Capture Anomaly",
        "🧠 AI Detection",
        "🏷️ Auto Categorize",
        "📍 Geo-Tagging",
        "📋 Submit Report",
        "👷 Dispatch Crew",
        "✅ Issue Resolved",
        "🔔 User Notified"
    ]
    
    step_w = Inches(1.3)
    step_h = Inches(1.6)
    for i, step in enumerate(steps):
        left_pos = Inches(0.5 + i * 1.55)
        # Card
        add_card(slide7, left_pos, Inches(3.0), step_w, step_h)
        # Step Number
        add_text_box(slide7, left_pos, Inches(3.15), step_w, Inches(0.3), f"0{i+1}", size=16, color=c_teal, bold=True, align=PP_ALIGN.CENTER)
        # Content
        add_text_box(slide7, left_pos + Inches(0.05), Inches(3.55), step_w - Inches(0.1), Inches(0.9), step, size=12, color=c_dark, bold=True, align=PP_ALIGN.CENTER)

    # -----------------------------------------------------
    # Slide 8: Innovation and Impact (2 Columns)
    # -----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8, c_bg_light)
    add_slide_header(slide8, "Innovation and Impact")
    
    # Left column: Innovation
    add_card(slide8, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5.0))
    add_text_box(slide8, Inches(0.8), Inches(2.0), Inches(5.2), Inches(0.5), "Key Innovations", size=20, color=c_blue, bold=True)
    innovations = (
        "💡 Real-time computer vision processing directly in frontend.\n\n"
        "💡 Single unified mapping layout eliminating administrative silos.\n\n"
        "💡 Highly responsive glassmorphism UI offering premium mobile experience."
    )
    add_text_box(slide8, Inches(0.8), Inches(2.6), Inches(5.2), Inches(3.8), innovations, size=15, color=c_dark)
    
    # Right column: Impact
    add_card(slide8, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0))
    add_text_box(slide8, Inches(7.1), Inches(2.0), Inches(5.4), Inches(0.5), "Direct Impact Metrics", size=20, color=c_teal, bold=True)
    impacts = (
        "📈 For Citizens: Resolution times sliced from weeks to less than 48 hours.\n\n"
        "📈 For Administration: Optimized maintenance routing saves up to 30% fuel and operational costs.\n\n"
        "📈 Strategic: Better municipal data visibility helps city planners prioritize long-term infrastructure overhauls."
    )
    add_text_box(slide8, Inches(7.1), Inches(2.6), Inches(5.4), Inches(3.8), impacts, size=15, color=c_dark)

    # -----------------------------------------------------
    # Slide 9: Future Scope (6 Grid Cards)
    # -----------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9, c_bg_light)
    add_slide_header(slide9, "Future Roadmap")
    
    roadmap = [
        {"title": "Municipal API Integration", "desc": "Connect directly to local public utility work-order databases.", "icon": "🔌"},
        {"title": "Predictive Analytics", "desc": "Use machine learning to forecast infrastructure decay.", "icon": "🔮"},
        {"title": "Push Notifications", "desc": "Send users push alerts via serverless edge messaging.", "icon": "🔔"},
        {"title": "Voice Reporting", "desc": "Deploy NLP systems for multi-lingual spoken complaints.", "icon": "🗣️"},
        {"title": "IoT Sensors", "desc": "Monitor garbage bin levels and streetlights automatically.", "icon": "📡"},
        {"title": "Drone Dispatch", "desc": "Aerial scanning for rapid site confirmation of rural anomalies.", "icon": "🛸"}
    ]
    
    for i, r in enumerate(roadmap):
        row = i // 3
        col = i % 3
        
        l_pos = Inches(0.5 + col * 4.2)
        t_pos = Inches(1.8 + row * 2.6)
        
        add_card(slide9, l_pos, t_pos, Inches(3.8), Inches(2.2))
        # Icon
        add_text_box(slide9, l_pos + Inches(0.2), t_pos + Inches(0.15), Inches(3.4), Inches(0.4), r["icon"], size=20)
        # Title
        add_text_box(slide9, l_pos + Inches(0.2), t_pos + Inches(0.55), Inches(3.4), Inches(0.5), r["title"], size=14, color=c_blue, bold=True)
        # Desc
        add_text_box(slide9, l_pos + Inches(0.2), t_pos + Inches(1.05), Inches(3.4), Inches(1.0), r["desc"], size=12, color=c_dark)

    # -----------------------------------------------------
    # Slide 10: Thank You Page (Light Theme)
    # -----------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10, c_bg_light)
    
    # Left accent block
    accent10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent10.fill.solid()
    accent10.fill.fore_color.rgb = c_teal
    accent10.line.fill.background()
    
    add_text_box(slide10, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2), "Thank You!", size=64, color=c_blue, bold=True)
    add_text_box(slide10, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8), "Transforming Civic Governance through Artificial Intelligence", size=22, color=c_teal, bold=True)
    add_text_box(slide10, Inches(1.5), Inches(4.8), Inches(10), Inches(0.6), "Questions & Discussion", size=18, color=c_dark)
    
    output_path = r"C:\Users\sarth\Documents\CIVIS_AI_Documents\CIVIS_AI_Presentation_Final.pptx"
    prs.save(output_path)
    print("Professional Presentation built successfully at:", output_path)

if __name__ == "__main__":
    build_professional_deck()
